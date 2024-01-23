import openpyxl
import re 
from pptx import Presentation
from pptx.util import Pt, Inches
from datetime import datetime

# Load the existing PowerPoint presentation template
presentation_path = 'QuadChartTemplate.pptx'
presentation = Presentation(presentation_path)

# Access the slide - usually the first slide
slide = presentation.slides[0]

# Set our placeholder indexes here, custom placeholders in PowerPoint start at 10
title_index = 10
upper_left_quad = 11
upper_right_quad = 12
poc_index = 13
lower_right_quad = 14
table_index = 15

# Clean the Excel file from any blanks and replace with unicode space
excel_file_path = 'UFRDATA-org.xlsx'
workbook = openpyxl.load_workbook(excel_file_path)

# Select the active sheet or a specific sheet by name
sheet = workbook.active

# Iterate through the rows and columns of the worksheet
for row in sheet.iter_rows():
    for cell in row:
        if cell.value is None:
            cell.value = '\u0020'  # Replace with unicode space

# Save the modified workbook
workbook.save('cleaned_ufrdata.xlsx')

# Close the workbook
workbook.close()

# Load the cleaned UFR Excel file
excel_file_path = 'cleaned_ufrdata.xlsx'
workbook = openpyxl.load_workbook(excel_file_path)

# Select the active sheet or a specific sheet by name
sheet = workbook.active

# Get the maximum row number
max_row = sheet.max_row

# Loop over rows
# Assuming data starts from the second row because of header (1-indexed)
for row_number in range(2, max_row + 1):  
    # Access data in each column of the current row
    submitting_org =        sheet.cell(row=row_number, column=1).value
    ufr_poc =               sheet.cell(row=row_number, column=2).value
    ufr_poc_phone =         sheet.cell(row=row_number, column=3).value
    ufr_sme =               sheet.cell(row=row_number, column=4).value
    ufr_sme_phone =         sheet.cell(row=row_number, column=5).value
    program =               sheet.cell(row=row_number, column=6).value
    funding_source =        sheet.cell(row=row_number, column=7).value
    ufr_title =             sheet.cell(row=row_number, column=8).value
    ufr_amount =            sheet.cell(row=row_number, column=9).value
    need_by =               sheet.cell(row=row_number, column=10).value
    comments =              sheet.cell(row=row_number, column=11).value
    description =           sheet.cell(row=row_number, column=12).value
    impact_if_not_funded =  sheet.cell(row=row_number, column=13).value
    manpower =              sheet.cell(row=row_number, column=14).value
    manpower_number =       sheet.cell(row=row_number, column=15).value
    support =               sheet.cell(row=row_number, column=16).value
    executable =            sheet.cell(row=row_number, column=17).value
    recurring =             sheet.cell(row=row_number, column=18).value
    incremental_funding =   sheet.cell(row=row_number, column=19).value
    fy24 =                  sheet.cell(row=row_number, column=20).value
    fy25 =                  sheet.cell(row=row_number, column=21).value
    fy26 =                  sheet.cell(row=row_number, column=22).value
    fy_total =              sheet.cell(row=row_number, column=23).value
    mitigation_action =     sheet.cell(row=row_number, column=24).value
    submitting_priority =   sheet.cell(row=row_number, column=25).value
    impact_score =          sheet.cell(row=row_number, column=26).value
    if_not_funded =         sheet.cell(row=row_number, column=27).value
    cdr_priority_loe =      sheet.cell(row=row_number, column=28).value
    mission_category =      sheet.cell(row=row_number, column=29).value
    contract =              sheet.cell(row=row_number, column=30).value
    cwg =                   sheet.cell(row=row_number, column=31).value
    pom_submission =        sheet.cell(row=row_number, column=32).value
    pom_year =              sheet.cell(row=row_number, column=33).value
    ddd =                   sheet.cell(row=row_number, column=34).value
    funding_category =      sheet.cell(row=row_number, column=35).value
    ba =                    sheet.cell(row=row_number, column=36).value
    pec =                   sheet.cell(row=row_number, column=37).value
    directorate_priority =  sheet.cell(row=row_number, column=38).value
    
    # Format our dates
    dateformat = "%d-%b-%y"
    ddd = ddd.strftime(dateformat)
    need_by = need_by.strftime(dateformat)
    
    # Add data to the slide
    # TODO: The add_paragraph() method adds a new line...How do we keep that from happening?
    # Upper title
    # TODO: I need to check to make sure if text length > 40, lower the font size
    # TODO: I need to figure out paragraph line height issues
    
    title_is_long = False
    
    if len(submitting_org) + len(ufr_title) > 40:
        upper_title = slide.shapes.placeholders[title_index]
        p = upper_title.text_frame.add_paragraph()
        r = p.add_run()
        r.text = submitting_org + " - " + ufr_title
        r = p.add_run()
        for paragraph in upper_title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
        upper_title.left = Inches(1.45)
        upper_title.top = Inches(0)
        upper_title.width = Inches(8.5)
        upper_title.height = Inches(1)
        
        # Flag title as long to reload PowerPoint template.  This will affect the other rows too.
        title_is_long = True
    else:
        upper_title = slide.shapes.placeholders[title_index]
        slide.shapes.placeholders[title_index].text = submitting_org
        slide.shapes.placeholders[title_index].text += " - " + ufr_title

    # Upper left quad
    ulq = slide.shapes.placeholders[upper_left_quad]
    p = ulq.text_frame.add_paragraph()
    r = p.add_run()
    r.text = "Description:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(description)
    r = p.add_run()

    r.text = "\nManpower Increase:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(manpower) + " "
    r = p.add_run()
    r.text = "How Many?"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(manpower_number)
    r = p.add_run()

    r.text = "\nCWG Approved:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(cwg) + " "
    r = p.add_run()
    r.text = "Support Agreement:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(support)
    r = p.add_run()

    # Upper Right quad
    urq = slide.shapes.placeholders[upper_right_quad]
    p = urq.text_frame.add_paragraph()
    r = p.add_run()
    r.text = "Impact if not Funded:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(impact_if_not_funded)
    r = p.add_run()
    
    r.text = "\nAction taken to mitigate this specific UFR:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(mitigation_action)
    r = p.add_run()
    
    # Lower left quad
    poc = slide.shapes.placeholders[poc_index]
    p = poc.text_frame.add_paragraph()
    r = p.add_run()
    r.text = "POC:"
    r.font.bold = True
    r = p.add_run()
    r.text = " " + str(ufr_poc) + ", " + str(submitting_org)
    r = p.add_run()

    # Make table
    table_placeholder = slide.placeholders[table_index]
    table = table_placeholder.table
    cell = table.cell(1, 0)
    cell.text = str(funding_category)
    cell = table.cell(1, 1)
    cell.text = "$" + str('{0:,.0f}'.format(int(fy24)/1000))
    cell = table.cell(1, 2)
    cell.text = "$" + str('{0:,.0f}'.format(int(fy25)/1000))
    cell = table.cell(1, 3)
    cell.text = "$" + str('{0:,.0f}'.format(int(fy26)/1000))
    cell = table.cell(1, 4)
    cell.text = "$" + str('{0:,.0f}'.format(int(fy_total)/1000))
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True

    
    # Format table cells from above
    i = 0
    while i < 5:
        cell = table.rows[1].cells[i]
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(13)
        i = i + 1
    
    # Lower Right quad
    lrq = slide.shapes.placeholders[lower_right_quad]
    p = lrq.text_frame.add_paragraph()
    r = p.add_run()
    r.text = "Funding Drop Dead Date (DDD):"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(ddd)
    r = p.add_run()
    
    r.text = "\nRequirement Identified in the POM:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(pom_submission)
    r = p.add_run()

    r.text = "\n  FY Submitted:"
    r.font.bold = True
    r = p.add_run()
    r.text = " " + str(need_by)
    r = p.add_run()
    
    r.text = "\nCategory:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(mission_category)
    r = p.add_run()
    
    r.text = "\nWhat CDR Priority/LOEs does this requirement support?:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(cdr_priority_loe)
    r = p.add_run()
    
    r.text = "\nCandidate for Incremental Funding:"
    r.font.bold = True
    r.font.underline = True
    r = p.add_run()
    r.text = " " + str(incremental_funding)
    r = p.add_run()

    # Save the updated presentation
    cleaned_filename = re.sub(r'[\\/*?:"<>|]'," ",ufr_title)
    updated_presentation_path = f'output\{cleaned_filename}.pptx'
    presentation.save(updated_presentation_path)
    print(f"Updated presentation saved to {updated_presentation_path}")
    
    # Clear the text frame or else it will keep appending data!
    upper_title.text_frame.clear()
    ulq.text_frame.clear()
    urq.text_frame.clear()
    poc.text_frame.clear()
    lrq.text_frame.clear()
    
    # Since we are modifying the title font and placeholder box we need to grab the template again.
    # TODO: If title is long, add to queue and then handle all rows with this issue instead of reloading.
    
    if title_is_long:
        # Reload the existing PowerPoint presentation template
        presentation_path = 'QuadChartTemplate.pptx'
        presentation = Presentation(presentation_path)

        # Access the slide - usually the first slide
        slide = presentation.slides[0]

# Close the workbook
workbook.close()

# Debugging Purposes: Used to find all placeholders and indexes
#for shape in slide.placeholders:
#    print('%d %s' % (shape.placeholder_format.idx, shape.name))